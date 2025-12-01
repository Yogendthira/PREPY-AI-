from flask import Flask, request, jsonify
from flask_cors import CORS
import os
from werkzeug.utils import secure_filename
import PyPDF2
from pptx import Presentation
import requests
import json
from dotenv import load_dotenv
from analysis import analyze_session

# Load environment variables
load_dotenv()

app = Flask(__name__)
CORS(app)

# Configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx'}
OLLAMA_API_URL = 'http://localhost:11434/api/chat'

print("✅ Using local Ollama model: phi3:3.8b")
print("📍 Ollama API: http://localhost:11434")

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size

# Create upload folder if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error extracting PDF: {e}")
    return text

def extract_text_from_ppt(file_path):
    """Extract text from PowerPoint file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting PPT: {e}")
    return text

def get_system_prompt(prep_type, difficulty, job_role=None):
    """Generate system prompt based on prep type, difficulty, and job role"""
    
    # Base prompts for each prep type
    base_prompts = {
   'interview': f"""You are an AI Interviewer{' for the position of ' + job_role if job_role else ''}.

⚠️ ABSOLUTE RULES - BREAKING THESE WILL FAIL THE TASK:
1. Your response MUST be EXACTLY ONE SHORT QUESTION
2. MAXIMUM 10 WORDS - Count them!
3. NO multiple questions - NO "and", NO commas separating questions
4. NO introductions, NO explanations, NO statements
5. Start directly with the question
6. End with a question mark

EXAMPLES OF CORRECT RESPONSES:
- "What technologies did you use?"
- "How does it work?"
- "What problem does this solve?"

EXAMPLES OF WRONG RESPONSES (TOO LONG):
- "Can you explain the architecture and how it scales?" (TWO questions!)
- "What specific features does your project incorporate to ensure accuracy?" (TOO LONG!)

Your goal: Ask ONE simple, direct question (max 10 words){' about the job role: ' + job_role if job_role else ''}.""",
        
        'hackathon': """You are a Hackathon Judge evaluating projects.

⚠️ ABSOLUTE RULES - BREAKING THESE WILL FAIL THE TASK:
1. Your response MUST be EXACTLY ONE SHORT QUESTION
2. MAXIMUM 10 WORDS - Count them!
3. NO multiple questions - NO "and", NO commas separating questions
4. NO introductions, NO explanations, NO statements
5. Start directly with the question
6. End with a question mark

EXAMPLES OF CORRECT RESPONSES:

- "what is the total build cost of this porject?"
- "How long did development take?"
- "How did your team collaborate?"

EXAMPLES OF WRONG RESPONSES (TOO LONG):
- "What features does it have and what are the limitations?" (TWO questions!)
- "How can you improve this project in future iterations?" (TOO LONG!)

Your goal: Ask ONE simple, direct question (max 10 words) about the project."""
    }
    
    # Difficulty modifiers - SIMPLIFIED
    difficulty_modifiers = {
        'superman': "Ask basic questions. Example: 'What does it do?'",
        'batman': "Ask practical questions. Example: 'How does it work?'",
        'hulk': "Ask technical questions. Example: 'What's the algorithm?'"
    }
    

    
    base = base_prompts.get(prep_type, base_prompts['interview'])
    modifier = difficulty_modifiers.get(difficulty, difficulty_modifiers['hulk'])
    
    return f"{base}\n\n{modifier}"


def call_ollama_api(messages, system_prompt):
    """Call local Ollama API with phi3 model"""
    try:
        # Ollama API endpoint
        url = 'http://localhost:11434/api/chat'
        
        # Format messages for Ollama
        formatted_messages = [{'role': 'system', 'content': system_prompt}] + messages
        
        payload = {
            'model': 'phi3:3.8b',
            'messages': formatted_messages,
            'stream': False,
            'options': {
                'num_predict': 60,  # Allow enough tokens to generate the full thought
                'temperature': 0.7,
                'top_p': 0.9
            }
        }
        
        response = requests.post(url, json=payload)
        response.raise_for_status()
        
        result = response.json()
        content = result['message']['content'].strip()
        
        # --- POST-PROCESSING: Extract ONLY the question ---
        # 1. If there's a question mark, take everything up to the first one
        if '?' in content:
            content = content.split('?')[0] + '?'
            
        # 2. If there are multiple sentences (split by . or !), take the last part
        # This removes "Intro text. Question?" -> "Question?"
        import re
        sentences = re.split(r'[.!]\s+', content)
        if sentences:
            content = sentences[-1].strip()
            
        return content
    except Exception as e:
        print(f"Ollama API Error: {e}")
        return f"Error communicating with AI: {str(e)}"

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Handle file upload and initial analysis"""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    prep_type = request.form.get('type', 'interview')
    difficulty = request.form.get('mode', 'hulk')
    job_role = request.form.get('job_role', '')
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Extract text based on file type
        if filename.endswith('.pdf'):
            extracted_text = extract_text_from_pdf(file_path)
        else:  # ppt or pptx
            extracted_text = extract_text_from_ppt(file_path)
        
        # Get system prompt
        system_prompt = get_system_prompt(prep_type, difficulty, job_role)
        
        # Store extracted text in session for later use
        # Use fixed welcome message instead of AI-generated one
        if prep_type == 'interview':
            ai_response = "Welcome to PREPY AI Interview. Give an Introduction about yourself."
        else:
            ai_response = "Welcome to PREPY AI Hackathon. Now Start with your project explanation."
        
        # Store the extracted text in session data for future questions
        session_data = {
            'prep_type': prep_type,
            'difficulty': difficulty,
            'system_prompt': system_prompt,
            'extracted_text': extracted_text
        }
        
        # Clean up uploaded file
        os.remove(file_path)
        
        return jsonify({
            'success': True,
            'message': ai_response,
            'extracted_text': extracted_text[:500],  # Send preview
            'session_data': session_data
        })
    
    return jsonify({'error': 'Invalid file type'}), 400

@app.route('/api/chat', methods=['POST'])
def chat():
    """Handle chat messages during evaluation"""
    data = request.json
    user_message = data.get('message', '')
    conversation_history = data.get('history', [])
    system_prompt = data.get('system_prompt', '')
    extracted_text = data.get('extracted_text', '')
    
    if not user_message:
        return jsonify({'error': 'No message provided'}), 400
    
    # Add user message to history
    conversation_history.append({'role': 'user', 'content': user_message})
    
    # Add context about uploaded content if this is early in conversation
    if len(conversation_history) <= 3 and extracted_text:
        context_message = f"Context from uploaded file:\n{extracted_text[:2000]}\n\nUser's response: {user_message}"
        conversation_history[-1]['content'] = context_message
    
    # Get AI response from Ollama
    ai_response = call_ollama_api(conversation_history, system_prompt)
    
    # Add AI response to history
    conversation_history.append({'role': 'assistant', 'content': ai_response})
    
    return jsonify({
        'success': True,
        'message': ai_response,
        'history': conversation_history
    })

@app.route('/api/analyze-session', methods=['POST'])
def analyze_session_endpoint():
    data = request.json
    history = data.get('history', [])
    job_role = data.get('job_role', 'Candidate')
    
    if not history:
        return jsonify({"success": False, "error": "No history provided"}), 400

    analysis_result = analyze_session(history, job_role)
    
    if analysis_result:
        return jsonify({"success": True, "data": analysis_result})
    else:
        # Fallback data if AI fails
        return jsonify({
            "success": False, 
            "error": "Analysis failed",
            "data": {
                "scores": {
                    "english": 0, "technical": 0, "communication": 0, 
                    "teamwork": 0, "soft_skills": 0, "project": 0, "overall": 0
                },
                "feedback": {
                    "strengths": "Analysis failed. Please try again.",
                    "improvements": "N/A",
                    "english_assessment": "N/A",
                    "recommendations": "N/A"
                }
            }
        })

@app.route('/api/save-recording', methods=['POST'])
def save_recording():
    """Save session recording to HACKATHONRECORDINGS folder"""
    try:
        if 'recording' not in request.files:
            return jsonify({'error': 'No recording provided'}), 400
        
        recording = request.files['recording']
        prep_type = request.form.get('type', 'unknown')
        mode = request.form.get('mode', 'unknown')
        
        # Create HACKATHONRECORDINGS folder if it doesn't exist
        recordings_folder = os.path.join(os.path.dirname(__file__), '..', 'HACKATHONRECORDINGS')
        os.makedirs(recordings_folder, exist_ok=True)
        
        # Save the recording
        filename = secure_filename(recording.filename)
        filepath = os.path.join(recordings_folder, filename)
        recording.save(filepath)
        
        print(f"✅ Recording saved: {filepath}")
        
        return jsonify({
            'success': True,
            'message': 'Recording saved successfully',
            'filepath': filepath,
            'filename': filename
        })
    except Exception as e:
        print(f"❌ Error saving recording: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({'status': 'healthy', 'message': 'Backend is running'})

if __name__ == '__main__':
    print("Starting Prepy AI Backend Server...")
    print("Server running on http://localhost:5000")
    app.run(debug=True, host='0.0.0.0', port=5000)
